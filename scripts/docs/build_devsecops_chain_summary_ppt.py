#!/usr/bin/env python3

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
OUT = ROOT / "docs" / "pptx" / "securerag-devsecops-chain-summary.pptx"

BG = RGBColor(247, 249, 252)
NAVY = RGBColor(47, 59, 82)
BLUE = RGBColor(43, 108, 176)
GREEN = RGBColor(34, 139, 85)
ORANGE = RGBColor(221, 107, 32)
AMBER = RGBColor(180, 83, 9)
TEXT = RGBColor(34, 34, 34)
MUTED = RGBColor(90, 100, 115)
WHITE = RGBColor(255, 255, 255)


def add_textbox(slide, left, top, width, height, text, size=20, bold=False,
                color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = "Aptos"
    font.size = Pt(size)
    font.bold = bold
    font.color.rgb = color
    p.alignment = align
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    return box


def add_card(slide, left, top, width, height, title, body, fill):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    tf = shape.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    r1.font.name = "Aptos Display"
    r1.font.size = Pt(20)
    r1.font.bold = True
    r1.font.color.rgb = WHITE
    p1.alignment = PP_ALIGN.LEFT
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = body
    r2.font.name = "Aptos"
    r2.font.size = Pt(12)
    r2.font.color.rgb = WHITE
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(10)
    tf.margin_bottom = Pt(10)
    return shape


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Slide 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = BG
add_textbox(slide, Inches(0.55), Inches(0.4), Inches(12.2), Inches(0.55),
            "SecureRAG Hub - Chaine DevSecOps et statuts honnetes", 24, True, NAVY)
add_textbox(slide, Inches(0.55), Inches(0.95), Inches(12.0), Inches(0.45),
            "Resume de la chaine end-to-end pour memoire et soutenance.", 12, False, MUTED)

card_w = Inches(3.0)
card_h = Inches(2.1)
add_card(slide, Inches(0.6), Inches(1.55), card_w, card_h,
         "TERMINÉ côté dépôt",
         "La chaine est conçue, automatisée, versionnée et documentée.",
         GREEN)
add_card(slide, Inches(3.85), Inches(1.55), card_w, card_h,
         "DÉPENDANT_DE_L_ENVIRONNEMENT",
         "Preuve runtime finale, release finale, Kyverno live, DB externe et Jenkins live.",
         BLUE)
add_card(slide, Inches(7.1), Inches(1.55), card_w, card_h,
         "PRÊT_NON_EXÉCUTÉ",
         "Modernisation secrets type SOPS, External Secrets Operator ou Vault.",
         ORANGE)
add_card(slide, Inches(10.35), Inches(1.55), Inches(2.35), card_h,
         "PARTIEL",
         "Démonstration live incomplète tant que le cluster cible n'est pas proprement rejoué.",
         AMBER)

add_textbox(slide, Inches(0.6), Inches(4.1), Inches(12.0), Inches(0.4),
            "Blocs principaux de la chaine", 17, True, NAVY)
bullets = [
    "1. Depot et gouvernance",
    "2. CI statique et controle de securite",
    "3. Build images et registry",
    "4. Supply chain et promotion digest",
    "5. Deploiement Kubernetes et runtime final",
    "6. Securite post-deploiement et Kyverno",
    "7. Donnees, secrets, Jenkins live, support pack final",
]
box = slide.shapes.add_textbox(Inches(0.8), Inches(4.55), Inches(11.8), Inches(2.2))
tf = box.text_frame
tf.word_wrap = True
tf.clear()
for i, line in enumerate(bullets):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = line
    p.font.name = "Aptos"
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT
    p.level = 0

# Slide 2
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = BG
add_textbox(slide, Inches(0.55), Inches(0.35), Inches(12.0), Inches(0.5),
            "Tableau de synthese de la chaine", 24, True, NAVY)

rows = [
    ("Code et gouvernance", "Git, Jenkinsfile, docs, runbooks", "TERMINÉ côté dépôt"),
    ("CI statique", "Lint, tests, Semgrep, Gitleaks, Trivy FS", "TERMINÉ côté dépôt"),
    ("Build et registry", "Docker, images OCI, localhost:5001", "TERMINÉ côté dépôt"),
    ("Supply chain", "Trivy image, Syft, Cosign sign/verify", "DÉPENDANT_DE_L_ENVIRONNEMENT"),
    ("Promotion release", "Digest-first, no-rebuild, attestation, provenance", "DÉPENDANT_DE_L_ENVIRONNEMENT"),
    ("Déploiement Kubernetes", "kind, overlays demo/production", "TERMINÉ côté dépôt"),
    ("Runtime final", "Pods récents, imageID, digests, HPA, healthchecks", "DÉPENDANT_DE_L_ENVIRONNEMENT"),
    ("Sécurité post-déploiement", "Hardening runtime, NetPol, SA, RBAC, PDB, probes", "DÉPENDANT_DE_L_ENVIRONNEMENT"),
    ("Kyverno Audit / Reports", "ClusterPolicies, PolicyReports", "DÉPENDANT_DE_L_ENVIRONNEMENT"),
    ("PostgreSQL externe", "Secret DB, backup, restore", "DÉPENDANT_DE_L_ENVIRONNEMENT"),
    ("Secrets modernes", "SOPS, ESO, Vault", "PRÊT_NON_EXÉCUTÉ"),
    ("Jenkins live", "Webhook, CI/CD reel, preuves", "PARTIEL"),
]
table = slide.shapes.add_table(len(rows) + 1, 3, Inches(0.45), Inches(1.0), Inches(12.35), Inches(5.85)).table
table.columns[0].width = Inches(2.8)
table.columns[1].width = Inches(5.9)
table.columns[2].width = Inches(3.65)
headers = ["Bloc", "Couverture", "Statut"]
for col, value in enumerate(headers):
    cell = table.cell(0, col)
    cell.text = value
    cell.fill.solid()
    cell.fill.fore_color.rgb = NAVY
    for p in cell.text_frame.paragraphs:
        for r in p.runs:
            r.font.name = "Aptos"
            r.font.size = Pt(12)
            r.font.bold = True
            r.font.color.rgb = WHITE

status_colors = {
    "TERMINÉ côté dépôt": GREEN,
    "DÉPENDANT_DE_L_ENVIRONNEMENT": BLUE,
    "PRÊT_NON_EXÉCUTÉ": ORANGE,
    "PARTIEL": AMBER,
}
for idx, (bloc, couverture, statut) in enumerate(rows, start=1):
    table.cell(idx, 0).text = bloc
    table.cell(idx, 1).text = couverture
    table.cell(idx, 2).text = statut
    for col in range(3):
        cell = table.cell(idx, col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE if idx % 2 else RGBColor(240, 244, 248)
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.name = "Aptos"
                r.font.size = Pt(10.5)
                r.font.color.rgb = TEXT
    table.cell(idx, 2).fill.fore_color.rgb = status_colors[statut]
    for p in table.cell(idx, 2).text_frame.paragraphs:
        for r in p.runs:
            r.font.bold = True
            r.font.color.rgb = WHITE

# Slide 3
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = BG
add_textbox(slide, Inches(0.55), Inches(0.4), Inches(12.0), Inches(0.5),
            "Lecture honnete pour le memoire et la soutenance", 24, True, NAVY)

left = Inches(0.7)
top = Inches(1.2)
sections = [
    ("Ce qui est déjà fort", [
        "Chaîne DevSecOps conçue, automatisée, versionnée et documentée.",
        "Déploiement digest strict prêt côté dépôt.",
        "Hardening Kubernetes, Kyverno Audit, support pack et documentation alignés.",
    ], GREEN),
    ("Ce qui doit être rejoué", [
        "Preuve runtime finale des workloads réellement actifs.",
        "Rejoue finale Trivy, Syft, Cosign, promotion digest, attestation et provenance.",
        "Collecte finale HPA, metrics, PolicyReports et runtime imageID.",
    ], BLUE),
    ("Ce qui reste environnemental", [
        "Kyverno Enforce si la registry n'est pas joignable depuis les pods Kyverno.",
        "PostgreSQL externe, backup et restore réels.",
        "Preuve Jenkins live sur l'instance réelle.",
    ], AMBER),
]
for title, lines, color in sections:
    add_textbox(slide, left, top, Inches(5.8), Inches(0.35), title, 18, True, color)
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top + Inches(0.35), Inches(5.7), Inches(1.45))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = color
    tf = box.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Aptos"
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT
        p.level = 0
    top += Inches(2.0)

add_textbox(slide, Inches(6.7), Inches(1.2), Inches(5.8), Inches(0.35),
            "Phrase de synthese", 18, True, NAVY)
summary = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.7), Inches(1.58), Inches(5.7), Inches(3.95))
summary.fill.solid()
summary.fill.fore_color.rgb = RGBColor(232, 240, 254)
summary.line.color.rgb = BLUE
tf = summary.text_frame
tf.clear()
paragraphs = [
    "SecureRAG Hub dispose d'une chaine DevSecOps complete et defendable cote depot.",
    "Les elements encore ouverts ne sont plus des trous de conception majeurs mais des preuves live a rejouer proprement sur l'environnement cible.",
    "La lecture honnete distingue explicitement ce qui est termine, partiel, pret non execute et dependant de l'environnement.",
]
for i, text in enumerate(paragraphs):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = text
    p.font.name = "Aptos"
    p.font.size = Pt(16 if i == 0 else 13)
    p.font.bold = i == 0
    p.font.color.rgb = TEXT

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(OUT)
