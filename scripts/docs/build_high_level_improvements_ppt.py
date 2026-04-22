from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("docs/pptx/securerag-high-level-improvements.pptx")


BG = RGBColor(245, 247, 250)
TEXT = RGBColor(28, 36, 52)
MUTED = RGBColor(96, 108, 128)
GREEN = RGBColor(29, 140, 87)
ORANGE = RGBColor(214, 135, 33)
BLUE = RGBColor(44, 102, 214)
RED = RGBColor(184, 63, 63)
CARD = RGBColor(255, 255, 255)
BORDER = RGBColor(218, 223, 232)


def add_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_title(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.9))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = TEXT
    if subtitle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = "Aptos"
        r2.font.size = Pt(11)
        r2.font.color.rgb = MUTED


def add_card(slide, x, y, w, h, title, body, color):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD
    shape.line.color.rgb = BORDER
    stripe = slide.shapes.add_shape(1, x, y, Inches(0.12), h)
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = color
    stripe.line.color.rgb = color

    tx = slide.shapes.add_textbox(x + Inches(0.22), y + Inches(0.12), w - Inches(0.3), h - Inches(0.2))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos"
    r.font.size = Pt(15)
    r.font.bold = True
    r.font.color.rgb = color
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = body
    r2.font.name = "Aptos"
    r2.font.size = Pt(11)
    r2.font.color.rgb = TEXT


def add_bullets(slide, x, y, w, h, items, font_size=16):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.level = 0
        p.bullet = True
        r = p.add_run()
        r.text = item
        r.font.name = "Aptos"
        r.font.size = Pt(font_size)
        r.font.color.rgb = TEXT


def slide_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(
        slide,
        "SecureRAG Hub - Ameliorations haut niveau",
        "Lecture honnete: renforcer l'environnement runtime, pas reecrire ce qui est deja solide cote depot.",
    )
    add_card(
        slide,
        Inches(0.7),
        Inches(1.5),
        Inches(5.8),
        Inches(1.4),
        "TERMINE cote depot",
        "La chaine est concue, automatisee, versionnee et documentee.",
        GREEN,
    )
    add_card(
        slide,
        Inches(6.7),
        Inches(1.5),
        Inches(5.8),
        Inches(1.4),
        "DEPENDANT_DE_L_ENVIRONNEMENT",
        "Runtime final, release finale, Kyverno live, DB externe et Jenkins live.",
        ORANGE,
    )
    add_card(
        slide,
        Inches(0.7),
        Inches(3.1),
        Inches(5.8),
        Inches(1.4),
        "PRET_NON_EXECUTE",
        "Modernisation secrets type SOPS, ESO ou Vault.",
        BLUE,
    )
    add_card(
        slide,
        Inches(6.7),
        Inches(3.1),
        Inches(5.8),
        Inches(1.4),
        "PARTIEL",
        "Demonstration live tant que le cluster cible n'est pas proprement rejoue.",
        RED,
    )


def slide_roadmap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Roadmap priorisee P0 / P1 / P2")
    add_card(
        slide,
        Inches(0.6),
        Inches(1.35),
        Inches(3.9),
        Inches(4.8),
        "P0 - Fermer les blocages",
        "Registry joignable par le cluster\nCluster cible sain et rejoue\nRelease finale par digest\nPostgreSQL externe avec backup et restore",
        RED,
    )
    add_card(
        slide,
        Inches(4.7),
        Inches(1.35),
        Inches(3.9),
        Inches(4.8),
        "P1 - Moderniser l'exploitation",
        "Secrets modernes\nPrometheus, Grafana, Alertmanager\nDetection runtime Falco ou Tetragon\nProgressive delivery",
        ORANGE,
    )
    add_card(
        slide,
        Inches(8.8),
        Inches(1.35),
        Inches(3.9),
        Inches(4.8),
        "P2 - Industrialiser la plateforme",
        "GitOps\nMapping conformite\nSLO, error budgets, incident drills\nGouvernance multi-environnements",
        BLUE,
    )


def slide_swot(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "SWOT haut niveau")
    add_card(slide, Inches(0.6), Inches(1.4), Inches(5.8), Inches(2.0), "Strengths", "Socle DevSecOps solide, release immuable, hardening Kubernetes, evidence model et documentation mature.", GREEN)
    add_card(slide, Inches(6.7), Inches(1.4), Inches(5.8), Inches(2.0), "Weaknesses", "Preuves runtime finales encore dependantes de l'environnement, registry loopback, DB externe non rejouee, Jenkins live partiel.", ORANGE)
    add_card(slide, Inches(0.6), Inches(3.7), Inches(5.8), Inches(2.0), "Opportunities", "Registry cloud, secrets modernes, observabilite SRE, GitOps, detection runtime, meilleure conformite.", BLUE)
    add_card(slide, Inches(6.7), Inches(3.7), Inches(5.8), Inches(2.0), "Threats", "Surpromesse production, environnement demo fragile, admission control limite par localhost, dette de gouvernance des secrets.", RED)


def slide_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Message de soutenance")
    add_bullets(
        slide,
        Inches(0.9),
        Inches(1.5),
        Inches(11.5),
        Inches(4.8),
        [
            "Le projet est deja fort la ou beaucoup de PFE restent superficiels: automatisation, supply chain, hardening, preuves et documentation.",
            "Les prochaines ameliorations ne sont pas des rustines de code; elles ciblent la credibilite runtime, la gouvernance des secrets, l'observabilite et la resilience des donnees.",
            "Le bon message n'est donc pas 'il manque tout', mais 'la base est solide et les prochains gains sont clairement identifies, priorises et operables'.",
        ],
        font_size=18,
    )
    footer = slide.shapes.add_textbox(Inches(0.9), Inches(6.5), Inches(11), Inches(0.5))
    p = footer.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "Priorites recommandees: registry reelle, cluster sain, PostgreSQL externe, secrets modernes, observabilite complete."
    r.font.name = "Aptos"
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = TEXT


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_overview(prs)
    slide_roadmap(prs)
    slide_swot(prs)
    slide_closing(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT.resolve())


if __name__ == "__main__":
    main()
